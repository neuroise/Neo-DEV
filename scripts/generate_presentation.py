#!/usr/bin/env python3
"""
Generate PowerPoint presentation for NO NOISE meeting.

Output: deliverables/NEUROISE_Presentazione_Marzo2026.pptx
"""

import json
from pathlib import Path
from pptx import Presentation
from pptx.util import Inches, Pt, Emu
from pptx.dml.color import RGBColor
from pptx.enum.text import PP_ALIGN, MSO_ANCHOR
from pptx.enum.shapes import MSO_SHAPE

BASE_DIR = Path(__file__).parent.parent
OUT_DIR = BASE_DIR / "deliverables"
OUT_DIR.mkdir(exist_ok=True)

# ── Colors ───────────────────────────────────────────────────────────────────
DARK_NAVY = RGBColor(0x1A, 0x1A, 0x2E)
MEDIUM_BLUE = RGBColor(0x2D, 0x3A, 0x5C)
ACCENT_BLUE = RGBColor(0x4A, 0x6F, 0xA5)
LIGHT_GRAY = RGBColor(0xF0, 0xF0, 0xF5)
WHITE = RGBColor(0xFF, 0xFF, 0xFF)
GREEN = RGBColor(0x27, 0xAE, 0x60)
YELLOW = RGBColor(0xF3, 0x9C, 0x12)
RED = RGBColor(0xE7, 0x4C, 0x3C)
TEXT_DARK = RGBColor(0x2C, 0x3E, 0x50)
TEXT_LIGHT = RGBColor(0x7F, 0x8C, 0x8D)

# Slide dimensions (16:9)
SLIDE_WIDTH = Inches(13.333)
SLIDE_HEIGHT = Inches(7.5)

# ── Data loading ─────────────────────────────────────────────────────────────

def load_summary(name):
    path = BASE_DIR / "data" / "experiments" / name / "summary.json"
    if path.exists():
        with open(path) as f:
            return json.load(f)
    return None

BASELINE = load_summary("baseline_30_llama70b_v3")
QWEN = load_summary("baseline_30_qwen32b")


# ── Helpers ──────────────────────────────────────────────────────────────────

def add_background(slide, color=DARK_NAVY):
    """Set slide background color."""
    background = slide.background
    fill = background.fill
    fill.solid()
    fill.fore_color.rgb = color


def add_shape_bg(slide, left, top, width, height, color, alpha=None):
    """Add a colored rectangle shape."""
    shape = slide.shapes.add_shape(MSO_SHAPE.RECTANGLE, left, top, width, height)
    shape.fill.solid()
    shape.fill.fore_color.rgb = color
    shape.line.fill.background()
    return shape


def add_text_box(slide, left, top, width, height, text, font_size=18,
                 color=WHITE, bold=False, alignment=PP_ALIGN.LEFT, font_name="Calibri"):
    """Add a text box with specified formatting."""
    txBox = slide.shapes.add_textbox(left, top, width, height)
    tf = txBox.text_frame
    tf.word_wrap = True
    p = tf.paragraphs[0]
    p.text = text
    p.font.size = Pt(font_size)
    p.font.color.rgb = color
    p.font.bold = bold
    p.font.name = font_name
    p.alignment = alignment
    return txBox


def add_bullet_list(slide, left, top, width, height, items, font_size=16,
                    color=WHITE, spacing=Pt(6)):
    """Add a bulleted list."""
    txBox = slide.shapes.add_textbox(left, top, width, height)
    tf = txBox.text_frame
    tf.word_wrap = True

    for i, item in enumerate(items):
        if i == 0:
            p = tf.paragraphs[0]
        else:
            p = tf.add_paragraph()
        p.text = item
        p.font.size = Pt(font_size)
        p.font.color.rgb = color
        p.font.name = "Calibri"
        p.space_after = spacing
        p.level = 0
    return txBox


def add_table_to_slide(slide, left, top, width, height, headers, rows,
                       header_color=MEDIUM_BLUE, row_colors=None):
    """Add a formatted table to a slide."""
    n_rows = len(rows) + 1
    n_cols = len(headers)
    table_shape = slide.shapes.add_table(n_rows, n_cols, left, top, width, height)
    table = table_shape.table

    # Header row
    for i, h in enumerate(headers):
        cell = table.cell(0, i)
        cell.text = h
        for p in cell.text_frame.paragraphs:
            p.font.size = Pt(12)
            p.font.bold = True
            p.font.color.rgb = WHITE
            p.font.name = "Calibri"
            p.alignment = PP_ALIGN.CENTER
        cell.fill.solid()
        cell.fill.fore_color.rgb = header_color

    # Data rows
    for r_idx, row in enumerate(rows):
        for c_idx, val in enumerate(row):
            cell = table.cell(r_idx + 1, c_idx)
            cell.text = str(val)
            for p in cell.text_frame.paragraphs:
                p.font.size = Pt(11)
                p.font.color.rgb = TEXT_DARK
                p.font.name = "Calibri"
                p.alignment = PP_ALIGN.CENTER
            if row_colors and c_idx < len(row_colors.get(r_idx, [])):
                rc = row_colors[r_idx][c_idx]
                if rc:
                    cell.fill.solid()
                    cell.fill.fore_color.rgb = rc

    return table


# ── Slides ───────────────────────────────────────────────────────────────────

def slide_title(prs):
    """Slide 1: Title"""
    slide = prs.slides.add_slide(prs.slide_layouts[6])  # Blank
    add_background(slide, DARK_NAVY)

    add_text_box(slide, Inches(1), Inches(1.5), Inches(11), Inches(1.5),
                 "NEURØISE Playground", font_size=44, bold=True, alignment=PP_ALIGN.CENTER)
    add_text_box(slide, Inches(1), Inches(3), Inches(11), Inches(1),
                 "Stato Avanzamento — Marzo 2026", font_size=28,
                 color=ACCENT_BLUE, alignment=PP_ALIGN.CENTER)
    add_text_box(slide, Inches(1), Inches(5), Inches(11), Inches(0.5),
                 "Collaborazione di Ricerca — Università di Pisa / NO NOISE S.r.l.",
                 font_size=16, color=TEXT_LIGHT, alignment=PP_ALIGN.CENTER)
    add_text_box(slide, Inches(1), Inches(5.8), Inches(11), Inches(0.5),
                 "Prof. Tiberio Uricchio",
                 font_size=14, color=TEXT_LIGHT, alignment=PP_ALIGN.CENTER)


def slide_agenda(prs):
    """Slide 2: Agenda"""
    slide = prs.slides.add_slide(prs.slide_layouts[6])
    add_background(slide, WHITE)

    add_text_box(slide, Inches(0.8), Inches(0.4), Inches(11), Inches(0.8),
                 "Agenda", font_size=36, color=DARK_NAVY, bold=True)

    items = [
        "1.  Contesto e Obiettivi del Progetto",
        "2.  Architettura del Sistema",
        "3.  Risultati Sperimentali",
        "4.  Servizio di Generazione Video",
        "5.  Stato Deliverables Contrattuali",
        "6.  Opportunità di Pubblicazione",
        "7.  Prossimi Passi",
    ]
    add_bullet_list(slide, Inches(1.5), Inches(1.8), Inches(10), Inches(5),
                    items, font_size=22, color=TEXT_DARK, spacing=Pt(14))


def slide_contesto(prs):
    """Slide 3: Project Context"""
    slide = prs.slides.add_slide(prs.slide_layouts[6])
    add_background(slide, WHITE)

    add_text_box(slide, Inches(0.8), Inches(0.4), Inches(11), Inches(0.8),
                 "Contesto del Progetto", font_size=36, color=DARK_NAVY, bold=True)

    left_items = [
        "Contratto: Collaborazione di Ricerca (19 Dic 2025)",
        "Durata: 3 mesi (scadenza ~19 Marzo 2026)",
        "Oggetto: MVP motore di storytelling intelligente",
        "Ambito: Esperienze personalizzate luxury / yachting",
    ]
    add_bullet_list(slide, Inches(0.8), Inches(1.6), Inches(5.5), Inches(3),
                    left_items, font_size=16, color=TEXT_DARK)

    right_items = [
        "Framework v2 → Cognitive Sandwich model",
        "Profilazione psicologica → Generazione video",
        "3 archetipi: Sage / Rebel / Lover",
        "30 profili ufficiali di test",
    ]
    add_bullet_list(slide, Inches(7), Inches(1.6), Inches(5.5), Inches(3),
                    right_items, font_size=16, color=TEXT_DARK)

    # Timeline bar
    add_shape_bg(slide, Inches(0.8), Inches(5.3), Inches(11.5), Inches(0.05), ACCENT_BLUE)
    months = [("Dic–Gen", "Analisi"), ("Febbraio", "Implementazione"), ("Marzo", "Esperimenti")]
    for i, (period, activity) in enumerate(months):
        x = Inches(1.5 + i * 3.8)
        add_text_box(slide, x, Inches(5.0), Inches(2.5), Inches(0.3),
                     period, font_size=12, color=ACCENT_BLUE, bold=True, alignment=PP_ALIGN.CENTER)
        add_text_box(slide, x, Inches(5.4), Inches(2.5), Inches(0.3),
                     activity, font_size=11, color=TEXT_LIGHT, alignment=PP_ALIGN.CENTER)


def slide_architecture(prs):
    """Slide 4: Cognitive Sandwich Architecture"""
    slide = prs.slides.add_slide(prs.slide_layouts[6])
    add_background(slide, WHITE)

    add_text_box(slide, Inches(0.8), Inches(0.4), Inches(11), Inches(0.8),
                 "Architettura — Cognitive Sandwich", font_size=36, color=DARK_NAVY, bold=True)

    layers_fixed = [
        ("INPUT LAYER", "Profilo utente (archetipo + musica + thread)", Inches(1.6), RGBColor(0x34, 0x98, 0xDB)),
        ("DECISION LAYER", "Director LLM → video triptych + OST prompt", Inches(2.8), ACCENT_BLUE),
        ("CONTROL LAYER", "PolicyGate → 8 regole, RED/YELLOW/GREEN", Inches(4.0), MEDIUM_BLUE),
        ("PRODUCTION LAYER", "Video Gen (Wan2.2) + Music Gen", Inches(5.2), DARK_NAVY),
    ]

    for label, desc, top, color in layers_fixed:
        shape = add_shape_bg(slide, Inches(1.5), top, Inches(10), Inches(0.9), color)
        add_text_box(slide, Inches(2), top + Inches(0.05), Inches(3), Inches(0.4),
                     label, font_size=16, bold=True, color=WHITE)
        add_text_box(slide, Inches(5.5), top + Inches(0.05), Inches(5.5), Inches(0.4),
                     desc, font_size=14, color=WHITE)

    # Arrows between layers
    for y in [Inches(2.55), Inches(3.75), Inches(4.95)]:
        add_text_box(slide, Inches(6), y, Inches(1), Inches(0.3),
                     "▼", font_size=18, color=TEXT_LIGHT, alignment=PP_ALIGN.CENTER)


def slide_pipeline(prs):
    """Slide 5: Playground Pipeline"""
    slide = prs.slides.add_slide(prs.slide_layouts[6])
    add_background(slide, WHITE)

    add_text_box(slide, Inches(0.8), Inches(0.4), Inches(11), Inches(0.8),
                 "Pipeline del Playground", font_size=36, color=DARK_NAVY, bold=True)

    steps = [
        ("Profile JSON", "30 profili\n3 archetipi"),
        ("Director LLM", "LLaMA 3.3:70b\nQwen3:32b"),
        ("PolicyGate", "8 regole\nRED/YELLOW/GREEN"),
        ("Metriche", "13 dimensioni\nautomatic"),
        ("Video Gen", "Wan2.2\nTurboWan"),
    ]

    for i, (title, desc) in enumerate(steps):
        x = Inches(0.5 + i * 2.5)
        # Box
        shape = add_shape_bg(slide, x, Inches(2.5), Inches(2.0), Inches(2.5), ACCENT_BLUE)
        add_text_box(slide, x + Inches(0.1), Inches(2.7), Inches(1.8), Inches(0.5),
                     title, font_size=16, bold=True, color=WHITE, alignment=PP_ALIGN.CENTER)
        add_text_box(slide, x + Inches(0.1), Inches(3.4), Inches(1.8), Inches(1.2),
                     desc, font_size=12, color=RGBColor(0xBD, 0xC3, 0xC7), alignment=PP_ALIGN.CENTER)

        # Arrow
        if i < len(steps) - 1:
            add_text_box(slide, x + Inches(2.0), Inches(3.2), Inches(0.5), Inches(0.5),
                         "→", font_size=24, color=TEXT_LIGHT, alignment=PP_ALIGN.CENTER)


def slide_director(prs):
    """Slide 6: Director LLM"""
    slide = prs.slides.add_slide(prs.slide_layouts[6])
    add_background(slide, WHITE)

    add_text_box(slide, Inches(0.8), Inches(0.4), Inches(11), Inches(0.8),
                 "Director LLM — Creative Agent", font_size=36, color=DARK_NAVY, bold=True)

    # Three archetype columns
    archetypes = [
        ("SAGE", "Il Saggio", "Contemplativo, minimale\n60-80 BPM\nAmbient, modern classical", GREEN),
        ("REBEL", "Il Ribelle", "Dinamico, audace\n120-140 BPM\nElectronic, breakbeat", RED),
        ("LOVER", "L'Amante", "Caldo, intimo\n70-90 BPM\nAcoustic, cinematic pop", RGBColor(0xE9, 0x1E, 0x63)),
    ]

    for i, (name, it_name, desc, color) in enumerate(archetypes):
        x = Inches(0.8 + i * 4)
        # Header
        shape = add_shape_bg(slide, x, Inches(1.5), Inches(3.5), Inches(0.7), color)
        add_text_box(slide, x + Inches(0.1), Inches(1.55), Inches(3.3), Inches(0.6),
                     f"{name} — {it_name}", font_size=18, bold=True, color=WHITE, alignment=PP_ALIGN.CENTER)
        # Description
        add_text_box(slide, x + Inches(0.2), Inches(2.4), Inches(3.1), Inches(2),
                     desc, font_size=14, color=TEXT_DARK, alignment=PP_ALIGN.CENTER)

    # Output section
    add_text_box(slide, Inches(0.8), Inches(4.8), Inches(11), Inches(0.5),
                 "Output: video_triptych (3 scene) + ost_prompt + metadata",
                 font_size=16, color=MEDIUM_BLUE, bold=True)

    add_text_box(slide, Inches(0.8), Inches(5.4), Inches(11), Inches(1),
                 "• 3 scene narrative: START → EVOLVE → END\n"
                 "• Prompt di produzione per text-to-video AI\n"
                 "• Vincolo: esclusivamente marino/costiero",
                 font_size=14, color=TEXT_DARK)


def slide_policygate(prs):
    """Slide 7: PolicyGate"""
    slide = prs.slides.add_slide(prs.slide_layouts[6])
    add_background(slide, WHITE)

    add_text_box(slide, Inches(0.8), Inches(0.4), Inches(11), Inches(0.8),
                 "PolicyGate — 8 Regole di Validazione", font_size=36, color=DARK_NAVY, bold=True)

    headers = ["Regola", "Tipo", "Descrizione"]
    rows = [
        ("R001", "RED", "Blacklist (urban, violence, forest...)"),
        ("R002", "YELLOW", "Warning terms (storm, person, logo...)"),
        ("R003", "YELLOW", "Marine vocabulary (≥5 termini)"),
        ("R004", "RED", "Structure (3 scene, ost_prompt)"),
        ("R005", "RED", "Scene sequence (start→evolve→end)"),
        ("R006", "YELLOW", "Archetype consistency (keyword density)"),
        ("R007", "YELLOW", "Prompt length (50-500 chars)"),
        ("R008", "RED/YEL", "BPM validation (presence + range)"),
    ]

    add_table_to_slide(slide, Inches(1), Inches(1.5), Inches(11), Inches(4),
                       headers, rows)

    # Compliance summary
    add_text_box(slide, Inches(1), Inches(5.8), Inches(11), Inches(1),
                 "Compliance baseline (LLaMA 70B):  Sage 80% GREEN  |  Rebel 20% GREEN  |  Lover 90% GREEN",
                 font_size=16, color=MEDIUM_BLUE, bold=True, alignment=PP_ALIGN.CENTER)


def slide_metrics(prs):
    """Slide 8: Metrics Framework"""
    slide = prs.slides.add_slide(prs.slide_layouts[6])
    add_background(slide, WHITE)

    add_text_box(slide, Inches(0.8), Inches(0.4), Inches(11), Inches(0.8),
                 "Framework Metriche — 13 Dimensioni", font_size=36, color=DARK_NAVY, bold=True)

    categories = [
        ("Strutturali", ["Schema Compliance", "Role Sequence", "Prompt Length", "Pacing"],
         RGBColor(0x34, 0x98, 0xDB)),
        ("Semantiche", ["Archetype Consistency", "Story Thread", "Lexical Fit", "Marine Vocab"],
         RGBColor(0x27, 0xAE, 0x60)),
        ("Qualitative", ["Red Flag Score", "Specificity", "Cross-Scene Coh.", "Narrative Coh."],
         RGBColor(0xE6, 0x7E, 0x22)),
        ("LLM Judge", ["Visual Clarity", "Arch. Alignment", "Emot. Resonance", "Marine Adh."],
         RGBColor(0x8E, 0x44, 0xAD)),
    ]

    for i, (cat_name, metrics, color) in enumerate(categories):
        x = Inches(0.5 + i * 3.1)
        # Category header
        shape = add_shape_bg(slide, x, Inches(1.5), Inches(2.8), Inches(0.6), color)
        add_text_box(slide, x + Inches(0.1), Inches(1.55), Inches(2.6), Inches(0.5),
                     cat_name, font_size=16, bold=True, color=WHITE, alignment=PP_ALIGN.CENTER)
        # Metric list
        add_bullet_list(slide, x + Inches(0.2), Inches(2.3), Inches(2.4), Inches(3),
                        metrics, font_size=13, color=TEXT_DARK, spacing=Pt(8))

    add_text_box(slide, Inches(0.8), Inches(5.8), Inches(11), Inches(0.5),
                 "Aggregate Score = media aritmetica di tutte le 13 metriche normalizzate [0, 1]",
                 font_size=14, color=TEXT_LIGHT, alignment=PP_ALIGN.CENTER)


def slide_baseline(prs):
    """Slide 9: Baseline Results"""
    slide = prs.slides.add_slide(prs.slide_layouts[6])
    add_background(slide, WHITE)

    add_text_box(slide, Inches(0.8), Inches(0.4), Inches(11), Inches(0.8),
                 "Risultati — Baseline (LLaMA 3.3:70b)", font_size=36, color=DARK_NAVY, bold=True)

    # Big number
    add_text_box(slide, Inches(1), Inches(1.5), Inches(5), Inches(1.5),
                 "0.775", font_size=72, color=ACCENT_BLUE, bold=True)
    add_text_box(slide, Inches(1), Inches(3), Inches(5), Inches(0.5),
                 "Aggregate Score (± 0.049)", font_size=20, color=TEXT_LIGHT)

    # Key stats
    stats = [
        "30/30 profili completati (100% success)",
        "30 profili × 13 metriche × 1 run",
        "Default prompt pack, temperature 0.7",
        "~4.9 tok/s su DGX Spark GB10",
    ]
    add_bullet_list(slide, Inches(1), Inches(3.8), Inches(5), Inches(3),
                    stats, font_size=16, color=TEXT_DARK)

    # Per-archetype boxes
    archetypes_data = [
        ("Sage", "0.805", "± 0.023", GREEN),
        ("Rebel", "0.723", "± 0.047", RED),
        ("Lover", "0.798", "± 0.018", RGBColor(0xE9, 0x1E, 0x63)),
    ]
    for i, (name, score, std, color) in enumerate(archetypes_data):
        x = Inches(7.5)
        y = Inches(1.8 + i * 1.6)
        shape = add_shape_bg(slide, x, y, Inches(4.5), Inches(1.2), color)
        add_text_box(slide, x + Inches(0.2), y + Inches(0.1), Inches(2), Inches(0.5),
                     name, font_size=20, bold=True, color=WHITE)
        add_text_box(slide, x + Inches(2.5), y + Inches(0.1), Inches(1.8), Inches(0.5),
                     score, font_size=28, bold=True, color=WHITE, alignment=PP_ALIGN.RIGHT)
        add_text_box(slide, x + Inches(2.5), y + Inches(0.6), Inches(1.8), Inches(0.4),
                     std, font_size=14, color=RGBColor(0xEC, 0xF0, 0xF1), alignment=PP_ALIGN.RIGHT)


def slide_ablation(prs):
    """Slide 10: Ablation Results"""
    slide = prs.slides.add_slide(prs.slide_layouts[6])
    add_background(slide, WHITE)

    add_text_box(slide, Inches(0.8), Inches(0.4), Inches(11), Inches(0.8),
                 "Risultati — Ablation Prompt Packs", font_size=36, color=DARK_NAVY, bold=True)

    headers = ["Metrica", "Default", "Concise", "Detailed"]
    rows = [
        ("Schema Compliance", "1.000", "1.000", "1.000"),
        ("Archetype Consistency", "0.789", "0.726", "0.770"),
        ("Story Thread", "0.822", "0.772", "0.711"),
        ("Red Flag Score", "0.755", "0.830", "0.835"),
        ("Cross-Scene Coherence", "0.588", "0.137", "0.140"),
        ("Prompt Specificity", "0.535", "0.542", "0.692"),
        ("LLM Judge", "0.770", "0.967", "0.958"),
        ("Aggregate", "0.775", "0.749", "0.744"),
    ]
    add_table_to_slide(slide, Inches(0.8), Inches(1.5), Inches(11.5), Inches(4.5),
                       headers, rows)

    add_text_box(slide, Inches(0.8), Inches(6.2), Inches(11), Inches(0.8),
                 "Default vince sull'aggregato. Concise sacrifica coerenza narrativa (5.3×). "
                 "Detailed migliora specificità ma aumenta varianza.",
                 font_size=14, color=TEXT_DARK)


def slide_crossmodel(prs):
    """Slide 11: Cross-Model Comparison"""
    slide = prs.slides.add_slide(prs.slide_layouts[6])
    add_background(slide, WHITE)

    add_text_box(slide, Inches(0.8), Inches(0.4), Inches(11), Inches(0.8),
                 "Risultati — Cross-Model Comparison", font_size=36, color=DARK_NAVY, bold=True)

    headers = ["Metrica", "LLaMA 3.3:70B", "Qwen3:32B", "p-value", "Effect (d)"]
    rows = [
        ("Schema", "1.000", "0.964", "0.157", "+0.39"),
        ("Archetype", "0.798", "0.770", "0.583", "+0.14"),
        ("Thread", "0.827", "0.720", "0.042*", "+0.45"),
        ("Red Flags", "0.748", "0.759", "0.714", "-0.09"),
        ("Coherence", "0.593", "0.102", "<.001***", "+5.46"),
        ("Specificity", "0.543", "0.731", "0.002**", "-0.67"),
        ("LLM Judge", "0.771", "0.996", "<.001***", "-2.37"),
        ("Aggregate", "0.777", "0.729", "<.001***", "+0.73"),
    ]
    add_table_to_slide(slide, Inches(0.5), Inches(1.5), Inches(12), Inches(4.5),
                       headers, rows)

    add_text_box(slide, Inches(0.8), Inches(6.2), Inches(11), Inches(0.8),
                 "LLaMA superiore nell'aggregato (p<.001). Differenza critica: "
                 "cross-scene coherence (d=5.46). Qwen eccelle in specificità e self-eval.",
                 font_size=14, color=TEXT_DARK)


def slide_archetype(prs):
    """Slide 12: Per-Archetype Analysis"""
    slide = prs.slides.add_slide(prs.slide_layouts[6])
    add_background(slide, WHITE)

    add_text_box(slide, Inches(0.8), Inches(0.4), Inches(11), Inches(0.8),
                 "Risultati — Analisi per Archetipo", font_size=36, color=DARK_NAVY, bold=True)

    headers = ["Metrica", "Sage", "Rebel", "Lover"]
    rows = [
        ("Archetype Consistency", "0.911", "0.589", "0.867"),
        ("Story Thread", "0.833", "0.733", "0.900"),
        ("Red Flag Score", "0.715", "0.775", "0.775"),
        ("Cross-Scene Coherence", "0.608", "0.504", "0.654"),
        ("Prompt Specificity", "0.752", "0.368", "0.485"),
        ("Pacing Progression", "0.997", "0.933", "0.924"),
        ("Aggregate", "0.805", "0.723", "0.798"),
    ]
    add_table_to_slide(slide, Inches(1), Inches(1.5), Inches(11), Inches(4),
                       headers, rows)

    # Key findings
    findings = [
        "Rebel: archetype consistency 0.589 (vs Sage 0.911) — 60% RED violations",
        "Sage: miglior aggregato (0.805) e specificità (0.752)",
        "Lover: miglior story thread (0.900) e coerenza (0.654)",
    ]
    add_bullet_list(slide, Inches(1), Inches(5.8), Inches(11), Inches(1.5),
                    findings, font_size=14, color=TEXT_DARK)


def slide_video(prs):
    """Slide 13: Video Generation"""
    slide = prs.slides.add_slide(prs.slide_layouts[6])
    add_background(slide, WHITE)

    add_text_box(slide, Inches(0.8), Inches(0.4), Inches(11), Inches(0.8),
                 "Generazione Video — Wan2.2 / TurboWan", font_size=36, color=DARK_NAVY, bold=True)

    left_items = [
        "Microservizio FastAPI containerizzato",
        "Due pipeline: Wan2.2 (qualità) / TurboWan (velocità)",
        "Hardware: DGX Spark GB10 (Blackwell GPU)",
        "121.7 GB unified memory",
        "Generazione batch per triptych (3 clip/profilo)",
        "Risoluzione configurabile, 5-10 sec per clip",
    ]
    add_bullet_list(slide, Inches(0.8), Inches(1.5), Inches(6), Inches(4),
                    left_items, font_size=16, color=TEXT_DARK)

    # Triptych visualization
    for i, scene in enumerate(["START", "EVOLVE", "END"]):
        x = Inches(8 + i * 1.6)
        shape = add_shape_bg(slide, x, Inches(2), Inches(1.3), Inches(2.5), MEDIUM_BLUE)
        add_text_box(slide, x + Inches(0.1), Inches(3), Inches(1.1), Inches(0.5),
                     scene, font_size=12, bold=True, color=WHITE, alignment=PP_ALIGN.CENTER)

    add_text_box(slide, Inches(8), Inches(4.8), Inches(5), Inches(0.5),
                 "Video Triptych", font_size=14, color=TEXT_LIGHT, alignment=PP_ALIGN.CENTER)


def slide_deliverables(prs):
    """Slide 14: Contract Deliverables Status"""
    slide = prs.slides.add_slide(prs.slide_layouts[6])
    add_background(slide, WHITE)

    add_text_box(slide, Inches(0.8), Inches(0.4), Inches(11), Inches(0.8),
                 "Deliverables Contrattuali — Stato", font_size=36, color=DARK_NAVY, bold=True)

    headers = ["Deliverable", "Descrizione", "Stato"]
    rows = [
        ("D1", "Documento di Analisi — Report architettura", "Completato ✓"),
        ("D2", "Software Playground — MVP funzionante (v0.2.0)", "Completato ✓"),
        ("D3", "Relazione Finale — Risultati e sviluppi futuri", "In corso (draft)"),
    ]
    add_table_to_slide(slide, Inches(1), Inches(1.8), Inches(11), Inches(2.5),
                       headers, rows)

    add_text_box(slide, Inches(1), Inches(5), Inches(11), Inches(1.5),
                 "• D1 e D2 completati nei termini\n"
                 "• D3 in fase di finalizzazione con risultati sperimentali completi\n"
                 "• Tutti i deliverables saranno consegnati entro la scadenza contrattuale",
                 font_size=16, color=TEXT_DARK)


def slide_paper(prs):
    """Slide 15: Publication Opportunity"""
    slide = prs.slides.add_slide(prs.slide_layouts[6])
    add_background(slide, DARK_NAVY)

    add_text_box(slide, Inches(0.8), Inches(0.4), Inches(11), Inches(0.8),
                 "Opportunità di Pubblicazione", font_size=36, color=WHITE, bold=True)

    add_text_box(slide, Inches(1), Inches(1.8), Inches(11), Inches(0.6),
                 "ACM Multimedia 2026", font_size=28, color=ACCENT_BLUE, bold=True)

    details = [
        "Conferenza di riferimento per multimedia intelligente (Rank A*)",
        "Deadline: ~Aprile 2026",
        "Topic: AI-driven personalized storytelling for luxury experiences",
        "Contributi: architettura, ablation study, cross-model evaluation",
        "",
        "Richiesta: autorizzazione alla pubblicazione congiunta",
        "Università di Pisa + NO NOISE S.r.l.",
    ]
    add_bullet_list(slide, Inches(1), Inches(2.8), Inches(11), Inches(4),
                    details, font_size=18, color=WHITE)


def slide_next_steps(prs):
    """Slide 16: Next Steps"""
    slide = prs.slides.add_slide(prs.slide_layouts[6])
    add_background(slide, WHITE)

    add_text_box(slide, Inches(0.8), Inches(0.4), Inches(11), Inches(0.8),
                 "Prossimi Passi", font_size=36, color=DARK_NAVY, bold=True)

    items = [
        ("Immediati (Marzo 2026)", [
            "Completamento e consegna D3 — Relazione Finale",
            "Esperimenti addizionali (cross-model ampliato)",
            "Finalizzazione framework valutazione umana",
        ]),
        ("Breve Termine (Aprile–Maggio 2026)", [
            "Draft paper ACM MM 2026 (previa autorizzazione)",
            "Campagna valutazione umana (5 dimensioni Likert)",
            "Ottimizzazione video generation pipeline",
        ]),
        ("Medio Termine", [
            "NeedsProfiler Agent (profilazione avanzata Big Five)",
            "Archivio narrativo (Summarizer + Memory Store)",
            "Deployment architecture per scenario on-yacht",
        ]),
    ]

    y = Inches(1.5)
    for section_title, section_items in items:
        add_text_box(slide, Inches(1), y, Inches(11), Inches(0.4),
                     section_title, font_size=18, color=ACCENT_BLUE, bold=True)
        y += Inches(0.5)
        add_bullet_list(slide, Inches(1.5), y, Inches(10), Inches(1.5),
                        section_items, font_size=15, color=TEXT_DARK, spacing=Pt(4))
        y += Inches(0.3) + Inches(len(section_items) * 0.35)


# ── Main ─────────────────────────────────────────────────────────────────────

def main():
    prs = Presentation()
    prs.slide_width = SLIDE_WIDTH
    prs.slide_height = SLIDE_HEIGHT

    slide_title(prs)       # 1
    slide_agenda(prs)      # 2
    slide_contesto(prs)    # 3
    slide_architecture(prs)  # 4
    slide_pipeline(prs)    # 5
    slide_director(prs)    # 6
    slide_policygate(prs)  # 7
    slide_metrics(prs)     # 8
    slide_baseline(prs)    # 9
    slide_ablation(prs)    # 10
    slide_crossmodel(prs)  # 11
    slide_archetype(prs)   # 12
    slide_video(prs)       # 13
    slide_deliverables(prs)  # 14
    slide_paper(prs)       # 15
    slide_next_steps(prs)  # 16

    out_path = OUT_DIR / "NEUROISE_Presentazione_Marzo2026.pptx"
    prs.save(str(out_path))
    print(f"✓ Presentation saved to: {out_path}")
    print(f"  Slides: {len(prs.slides)}")
    print(f"  Size: {out_path.stat().st_size / 1024:.1f} KB")


if __name__ == "__main__":
    main()
